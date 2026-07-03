"""独立测试混合导出功能，使用已有任务文件，不需要调用 LLM。

用法:
    .venv\\Scripts\\python.exe -m pytest tests\\test_hybrid_export.py -x -s
"""

import json
from pathlib import Path

import pytest
from pptx import Presentation

from app.core.config import Settings
from app.infrastructure.ppt_master.pptx_builder import PPTBuilder
from app.services.hybrid_pptx_exporter import HybridPptxExporter
from app.schemas.structured_generation import StructuredPageResult, GeneratedElement

# 任务目录
TASK_DIR = Path("mock_ftp/slides_gen_server/tasks/task_20260703160147_7358a679")
# 模板目录
TEMPLATE_DIR = Path("mock_ftp/slides_gen_server/templates/tpl_20260627132835_bce99ff8")
TEMPLATE_PPTX = TEMPLATE_DIR / "source" / "template.pptx"
# 模板规则（从 runtime 中取最新的）
TEMPLATE_RULES_PATH = Path("runtime/tasks/task_20260703152015_fbecc901/analysis/template_rules.json")

# 本地数据是否可用
_LOCAL_DATA_AVAILABLE = TEMPLATE_PPTX.exists() and TEMPLATE_RULES_PATH.exists() and TASK_DIR.exists()

_skip_if_no_local_data = pytest.mark.skipif(
    not _LOCAL_DATA_AVAILABLE,
    reason="本地任务数据不存在，跳过（CI 环境无此数据）",
)


def _make_settings(tmp_path: Path) -> Settings:
    """构造测试用 Settings。"""
    return Settings(
        app_name="test",
        app_env="test",
        api_prefix="/api/v1",
        runtime_dir=tmp_path / "runtime",
        mock_ftp_dir=tmp_path / "mock_ftp",
        default_template_file=tmp_path / "template.pptx",
        db_host="127.0.0.1",
        db_port=3306,
        db_user="root",
        db_password="",
        db_schema="test",
        ftp_host="",
        ftp_port=21,
        ftp_user="",
        ftp_password="",
        ftp_root_dir="/slides_gen_server",
        mock_ftp_enabled=True,
        default_template_id=None,
        ppt_master_scripts_dir=Path("app/vendor/ppt_master/scripts").resolve(),
        llm_base_url="https://api.deepseek.com",
        llm_model="deepseek-v4-flash",
        llm_timeout_seconds=120,
        max_llm_concurrency=8,
        llm_rate_limit_max_retries=3,
        llm_rate_limit_base_delay=1.0,
        llm_rate_limit_max_delay=60.0,
        svg_page_types="diagram",
    )


def _load_structured_results() -> dict[int, StructuredPageResult]:
    """从 structured_results 目录加载所有结构化生成结果。"""
    results_dir = TASK_DIR / "structured_results"
    pages: dict[int, StructuredPageResult] = {}
    for json_file in sorted(results_dir.glob("page_*.json")):
        data = json.loads(json_file.read_text(encoding="utf-8"))
        elements = []
        for elem_data in data.get("elements", []):
            elements.append(GeneratedElement(
                id=elem_data["id"],
                type=elem_data["type"],
                content=elem_data.get("content"),
                headers=elem_data.get("headers"),
                rows=elem_data.get("rows"),
            ))
        result = StructuredPageResult(
            page_no=data["page_no"],
            should_generate=data["should_generate"],
            skip_reason=data.get("skip_reason", ""),
            elements=elements,
        )
        pages[result.page_no] = result
    return pages


def _load_skipped_pages() -> set[int]:
    """从 page_plans.json 中提取跳过的页码。"""
    plans_path = TASK_DIR / "analysis" / "page_plans.json"
    if not plans_path.exists():
        return set()
    plans = json.loads(plans_path.read_text(encoding="utf-8"))
    skipped = set()
    for plan in plans:
        if not plan.get("should_generate", True):
            skipped.add(plan["page_no"])
    return skipped


def _load_svg_pages() -> dict[int, Path]:
    """从 svg_final 目录加载所有 SVG 文件。"""
    svg_dir = TASK_DIR / "svg_final"
    pages: dict[int, Path] = {}
    for svg_file in sorted(svg_dir.glob("slide_*.svg")):
        # slide_15.svg -> page 15
        page_no = int(svg_file.stem.split("_")[1])
        pages[page_no] = svg_file
    return pages


@_skip_if_no_local_data
def test_hybrid_export_with_real_data(tmp_path: Path):
    """使用真实任务数据测试混合导出。"""
    # 加载所有数据
    structured_pages = _load_structured_results()
    skipped_pages = _load_skipped_pages()
    svg_pages = _load_svg_pages()
    template_rules = json.loads(TEMPLATE_RULES_PATH.read_text(encoding="utf-8"))

    print(f"\n结构化页面: {sorted(structured_pages.keys())}")
    print(f"跳过页面: {sorted(skipped_pages)}")
    print(f"SVG 页面: {sorted(svg_pages.keys())}")
    print(f"模板规则页数: {len(template_rules.get('pages', []))}")

    # 构造导出器
    settings = _make_settings(tmp_path)
    exporter = HybridPptxExporter(settings=settings)

    output_path = tmp_path / "result.pptx"

    # 执行导出
    result_path = exporter.export(
        template_pptx_path=TEMPLATE_PPTX,
        template_rules=template_rules,
        svg_pages=svg_pages,
        structured_pages=structured_pages,
        skipped_pages=skipped_pages,
        output_path=output_path,
    )

    print(f"\n导出完成: {result_path}")
    print(f"文件大小: {result_path.stat().st_size} bytes")

    # 验证输出文件
    assert result_path.exists(), "输出文件不存在"
    assert result_path.stat().st_size > 0, "输出文件为空"

    # 打开验证 slide 数量
    pres = Presentation(result_path)
    print(f"输出 PPTX slide 数量: {len(pres.slides)}")
    # 预期: 48 页 - 跳过的页面数 + 拆页（如果有）
    total_input = 48
    skipped_count = len(skipped_pages)
    # should_generate=False 的结构化页面也算跳过
    for page_no, result in structured_pages.items():
        if not result.should_generate:
            skipped_count += 1
    expected_min = total_input - skipped_count
    print(f"预期最少 slide 数: {expected_min} (总{total_input} - 跳过{skipped_count})")
    assert len(pres.slides) >= expected_min, f"slide 数量 {len(pres.slides)} 少于预期 {expected_min}"


@_skip_if_no_local_data
def test_fill_single_slide_basic(tmp_path: Path):
    """测试单页结构化填充。"""
    template_rules = json.loads(TEMPLATE_RULES_PATH.read_text(encoding="utf-8"))
    structured_pages = _load_structured_results()

    # 找一个 should_generate=True 的页面
    test_page_no = None
    for page_no, result in structured_pages.items():
        if result.should_generate and result.elements:
            test_page_no = page_no
            break

    if test_page_no is None:
        pytest.skip("没有可用的结构化生成结果")

    print(f"\n测试页面: {test_page_no}")
    result = structured_pages[test_page_no]
    print(f"元素数量: {len(result.elements)}")

    # 获取该页的规则
    pages_list = template_rules.get("pages", [])
    page_rule = pages_list[test_page_no - 1] if test_page_no <= len(pages_list) else {}
    print(f"页面规则元素数: {len(page_rule.get('elements', []))}")

    # 填充单页
    pres = Presentation(TEMPLATE_PPTX)
    builder = PPTBuilder(TEMPLATE_PPTX)
    slide_index = test_page_no - 1
    occupied = builder.fill_single_slide(pres, slide_index, page_rule, result)

    print(f"占用 slide 索引: {occupied}")

    # 保存
    output = tmp_path / "single_slide.pptx"
    pres.save(output)
    print(f"保存到: {output} ({output.stat().st_size} bytes)")
    assert output.exists()
