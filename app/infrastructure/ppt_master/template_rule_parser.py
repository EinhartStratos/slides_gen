"""模板规则解析器。

从 pptx_gen/template/parser.py 移植，去掉图片/Mermaid 相关逻辑，
适配 app.schemas.structured_generation 中的数据模型。
解析 PPTX 模板，提取每页的元素规则（文本框位置、表格结构等）。
"""

from __future__ import annotations

from pathlib import Path
import json
import re
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation

from app.schemas.structured_generation import (
    BoundingBox,
    ElementRule,
    PageRule,
    SlideSize,
    TableSchema,
    TemplateMeta,
    TemplateRules,
    TitleRule,
)


class TemplateRuleParser:
    """解析 PPTX 模板，提取每页元素规则。"""

    NS = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    def __init__(self, template_path: str | Path) -> None:
        self.template_path = Path(template_path)

    def parse(self) -> TemplateRules:
        presentation = Presentation(self.template_path)
        with zipfile.ZipFile(self.template_path) as archive:
            slide_names = self._sorted_slide_names(archive)
            pages = [
                self._parse_slide(archive, slide_name, presentation.slides[index - 1], index)
                for index, slide_name in enumerate(slide_names, start=1)
            ]
        template = TemplateMeta(
            file_name=self.template_path.name,
            slide_count=len(pages),
            slide_size=SlideSize(cx=presentation.slide_width, cy=presentation.slide_height),
        )
        return TemplateRules(template=template, pages=pages)

    def save(self, output_path: str | Path) -> TemplateRules:
        """解析模板并保存规则 JSON 到指定路径。"""
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        rules = self.parse()
        output.write_text(
            json.dumps(rules.model_dump(mode="json", exclude_none=True), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return rules

    def _sorted_slide_names(self, archive: zipfile.ZipFile) -> list[str]:
        names = [name for name in archive.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)]
        return sorted(names, key=lambda item: int(re.search(r"slide(\d+)\.xml$", item).group(1)))

    def _parse_slide(self, archive: zipfile.ZipFile, slide_name: str, slide, page_no: int) -> PageRule:
        root = ET.fromstring(archive.read(slide_name))
        sp_tree = root.find(".//p:spTree", self.NS)
        if sp_tree is None:
            return PageRule(
                page_no=page_no,
                page_name=f"slide_{page_no}",
                page_purpose="empty",
                supports_mermaid=False,
                title=TitleRule(),
                elements=[],
            )

        title_info = TitleRule()
        elements: list[ElementRule] = []
        z_order = 0

        for child in list(sp_tree):
            tag = self._local_name(child.tag)
            if tag == "sp":
                parsed = self._parse_shape(child, page_no, z_order)
                z_order += 1
                if parsed is None:
                    continue
                if parsed.type == "title":
                    title_info = TitleRule(
                        text=parsed.default_text or "",
                        bbox=parsed.bbox,
                        style=parsed.style,
                    )
                elements.append(parsed)
            elif tag == "graphicFrame":
                parsed = self._parse_table(child, slide, page_no, z_order)
                z_order += 1
                if parsed is not None:
                    elements.append(parsed)
            else:
                z_order += 1

        page_name = title_info.text or f"slide_{page_no}"
        page_purpose = self._infer_page_purpose(page_name, elements)
        return PageRule(
            page_no=page_no,
            page_name=page_name,
            page_purpose=page_purpose,
            supports_mermaid=False,
            title=title_info,
            elements=elements,
        )

    def _parse_shape(self, node: ET.Element, page_no: int, z_order: int) -> ElementRule | None:
        c_nv_pr = node.find("./p:nvSpPr/p:cNvPr", self.NS)
        if c_nv_pr is None:
            return None
        shape_id = int(c_nv_pr.attrib.get("id", "0"))
        name = c_nv_pr.attrib.get("name", f"shape_{shape_id}")
        bbox = self._parse_bbox(node.find("./p:spPr/a:xfrm", self.NS))
        default_text = self._normalize_default_text(self._extract_text(node))
        placeholder = node.find("./p:nvSpPr/p:nvPr/p:ph", self.NS)
        is_title = placeholder is not None and placeholder.attrib.get("type") == "title"
        element_type = "title" if is_title else "text"
        is_instructional = False if is_title else self._is_instructional_text(name, default_text)
        role = "title" if is_title else self._infer_text_role(name, default_text, is_instructional)
        style = self._parse_text_style(node)
        content_requirement = self._infer_content_requirement(element_type, default_text, name, is_instructional)
        return ElementRule(
            id=f"{element_type}_{shape_id}",
            shape_id=shape_id,
            type=element_type,
            role=role,
            page_no=page_no,
            bbox=bbox,
            z_order=z_order,
            editable=True,
            default_text=default_text,
            style=style,
            content_requirement=content_requirement,
            fill_strategy="replace_text",
            is_instructional=is_instructional,
        )

    def _parse_table(self, node: ET.Element, slide, page_no: int, z_order: int) -> ElementRule | None:
        table = node.find(".//a:tbl", self.NS)
        c_nv_pr = node.find("./p:nvGraphicFramePr/p:cNvPr", self.NS)
        xfrm = node.find("./p:xfrm", self.NS)
        if table is None or c_nv_pr is None:
            return None
        shape_id = int(c_nv_pr.attrib.get("id", "0"))
        bbox = self._parse_bbox(xfrm)
        rows = table.findall("./a:tr", self.NS)
        cols = table.findall("./a:tblGrid/a:gridCol", self.NS)
        cell_texts: list[list[str]] = []
        for row in rows:
            values = []
            for cell in row.findall("./a:tc", self.NS):
                values.append(self._extract_text(cell))
            cell_texts.append(values)
        shape = next(
            (item for item in slide.shapes if item.shape_id == shape_id and getattr(item, "has_table", False)),
            None,
        )
        style = self._parse_table_style(shape) if shape is not None else self._default_table_style()
        return ElementRule(
            id=f"table_{shape_id}",
            shape_id=shape_id,
            type="table",
            role="table",
            page_no=page_no,
            bbox=bbox,
            z_order=z_order,
            editable=True,
            default_text=self._normalize_default_text("\n".join([" | ".join(row) for row in cell_texts])),
            style=style,
            content_requirement="按模板表格结构填写内容，输出行列清晰、不要超过模板容量。",
            fill_strategy="fill_table",
            table_schema=TableSchema(
                rows=len(rows),
                cols=len(cols),
                default_cells=cell_texts or None,
                max_rows=len(rows),
                max_cols=len(cols),
            ),
        )

    def _parse_bbox(self, xfrm: ET.Element | None) -> BoundingBox:
        if xfrm is None:
            return BoundingBox()
        off = xfrm.find("./a:off", self.NS)
        ext = xfrm.find("./a:ext", self.NS)
        return BoundingBox(
            x=int(off.attrib.get("x", "0")) if off is not None else 0,
            y=int(off.attrib.get("y", "0")) if off is not None else 0,
            w=int(ext.attrib.get("cx", "0")) if ext is not None else 0,
            h=int(ext.attrib.get("cy", "0")) if ext is not None else 0,
        )

    def _parse_text_style(self, node: ET.Element) -> dict[str, object] | None:
        body_pr = node.find("./p:txBody/a:bodyPr", self.NS)
        paragraph = node.find("./p:txBody/a:p", self.NS)
        first_run = node.find("./p:txBody/a:p/a:r/a:rPr", self.NS)
        style: dict[str, object] = {}
        if body_pr is not None:
            style["margins"] = {
                "left": int(body_pr.attrib.get("lIns", "0")),
                "top": int(body_pr.attrib.get("tIns", "0")),
                "right": int(body_pr.attrib.get("rIns", "0")),
                "bottom": int(body_pr.attrib.get("bIns", "0")),
            }
            if "anchor" in body_pr.attrib:
                style["vertical_anchor"] = body_pr.attrib["anchor"]
        if paragraph is not None:
            paragraph_pr = paragraph.find("./a:pPr", self.NS)
            if paragraph_pr is not None and "algn" in paragraph_pr.attrib:
                style["alignment"] = paragraph_pr.attrib["algn"]
        if first_run is not None:
            if "sz" in first_run.attrib:
                style["font_size"] = int(first_run.attrib["sz"])
            if "b" in first_run.attrib:
                style["bold"] = first_run.attrib["b"] == "1"
            if "i" in first_run.attrib:
                style["italic"] = first_run.attrib["i"] == "1"
            latin = first_run.find("./a:latin", self.NS)
            if latin is not None and "typeface" in latin.attrib:
                style["font_name"] = latin.attrib["typeface"]
            color = first_run.find("./a:solidFill/a:srgbClr", self.NS)
            scheme_color = first_run.find("./a:solidFill/a:schemeClr", self.NS)
            if color is not None:
                style["font_color"] = color.attrib.get("val", "")
            elif scheme_color is not None:
                style["font_color_scheme"] = scheme_color.attrib.get("val", "")
        return style or None

    def _parse_table_style(self, shape) -> dict[str, object]:
        default_style = self._default_table_style()
        table = shape.table
        header_cells = [table.cell(0, col_index) for col_index in range(len(table.columns))] if len(table.rows) else []
        body_cells = []
        fallback_body_cells = []
        for row_index in range(1, len(table.rows)):
            for col_index in range(len(table.columns)):
                cell = table.cell(row_index, col_index)
                if cell.text.strip():
                    fallback_body_cells.append(cell)
                if self._is_meaningful_body_cell_text(cell.text):
                    body_cells.append(cell)
        header_style = self._extract_table_cell_style(header_cells)
        body_style = self._extract_table_cell_style(body_cells) or self._extract_table_cell_style(fallback_body_cells)
        resolved_header = self._merge_table_text_style(default_style["header_style"], header_style)
        resolved_body = self._merge_table_text_style(default_style["body_style"], body_style)
        return {
            "font_name": resolved_body["font_name"],
            "font_size": resolved_body["font_size"],
            "bold": resolved_body.get("bold", False),
            "italic": resolved_body.get("italic", False),
            "alignment": resolved_body.get("alignment", "left"),
            "vertical_anchor": resolved_body.get("vertical_anchor", "middle"),
            "margins": dict(resolved_body.get("margins") or {}),
            "row_heights": [int(row.height) for row in table.rows],
            "column_widths": [int(column.width) for column in table.columns],
            "header_style": resolved_header,
            "body_style": resolved_body,
        }

    def _extract_table_cell_style(self, cells: list) -> dict[str, object] | None:
        for cell in cells:
            text_frame = cell.text_frame
            paragraph = next((item for item in text_frame.paragraphs if item.text.strip()), None)
            if paragraph is None and text_frame.paragraphs:
                paragraph = text_frame.paragraphs[0]
            if paragraph is None:
                continue
            run = next((item for item in paragraph.runs if item.text.strip()), None)
            if run is None and paragraph.runs:
                run = paragraph.runs[0]
            style: dict[str, object] = {
                "margins": {
                    "left": int(cell.margin_left),
                    "top": int(cell.margin_top),
                    "right": int(cell.margin_right),
                    "bottom": int(cell.margin_bottom),
                }
            }
            if paragraph.alignment is not None:
                style["alignment"] = self._normalize_alignment(paragraph.alignment)
            if cell.vertical_anchor is not None:
                style["vertical_anchor"] = self._normalize_vertical_anchor(cell.vertical_anchor)
            if run is not None:
                font = run.font
                if font.name:
                    style["font_name"] = font.name
                if font.size:
                    style["font_size"] = int(round(font.size.pt * 100))
                if font.bold is not None:
                    style["bold"] = font.bold
                if font.italic is not None:
                    style["italic"] = font.italic
            if style:
                return style
        return None

    def _default_table_style(self) -> dict[str, object]:
        header_style = {
            "font_name": "微软雅黑",
            "font_size": 1400,
            "bold": True,
            "italic": False,
            "alignment": "center",
            "vertical_anchor": "middle",
            "margins": {"left": 91440, "top": 45720, "right": 91440, "bottom": 45720},
        }
        body_style = {
            "font_name": "微软雅黑",
            "font_size": 1400,
            "bold": False,
            "italic": False,
            "alignment": "left",
            "vertical_anchor": "middle",
            "margins": {"left": 68580, "top": 0, "right": 68580, "bottom": 0},
        }
        return {
            "font_name": body_style["font_name"],
            "font_size": body_style["font_size"],
            "bold": body_style["bold"],
            "italic": body_style["italic"],
            "alignment": body_style["alignment"],
            "vertical_anchor": body_style["vertical_anchor"],
            "margins": dict(body_style["margins"]),
            "row_heights": [],
            "column_widths": [],
            "header_style": header_style,
            "body_style": body_style,
        }

    def _merge_table_text_style(self, base_style: dict[str, object], override_style: dict[str, object] | None) -> dict[str, object]:
        merged = dict(base_style)
        merged["margins"] = dict(base_style.get("margins") or {})
        if not override_style:
            return merged
        for key, value in override_style.items():
            if key == "margins":
                merged["margins"].update(value or {})
            else:
                merged[key] = value
        return merged

    def _is_meaningful_body_cell_text(self, text: str) -> bool:
        normalized = re.sub(r"[\s\d…·,.，。:：;；()（）/\\\-]+", "", text or "")
        return bool(normalized)

    def _normalize_alignment(self, alignment) -> str:
        name = getattr(alignment, "name", "")
        if not name:
            return "left"
        lowered = name.lower()
        mapping = {
            "left": "left", "center": "center", "right": "right",
            "justify": "justify", "distribute": "distribute",
            "thai_distribute": "distribute", "justify_low": "justify",
        }
        return mapping.get(lowered, "left")

    def _normalize_vertical_anchor(self, anchor) -> str:
        name = getattr(anchor, "name", "")
        if not name:
            return "middle"
        lowered = name.lower()
        mapping = {"top": "top", "middle": "middle", "bottom": "bottom"}
        return mapping.get(lowered, "middle")

    def _extract_text(self, node: ET.Element) -> str:
        texts = [text.text or "" for text in node.findall(".//a:t", self.NS)]
        return "".join(texts).strip()

    def _normalize_default_text(self, text: str) -> str | None:
        normalized = re.sub(r"\s+", " ", text).strip()
        return normalized or None

    def _infer_page_purpose(self, page_name: str, elements: list[ElementRule]) -> str:
        lowered = page_name.lower()
        if "封面" in page_name or "cover" in lowered:
            return "cover"
        if any(item.type == "table" for item in elements):
            return "table"
        return "text"

    def _infer_text_role(self, shape_name: str, default_text: str | None, is_instructional: bool) -> str:
        content = f"{shape_name} {default_text or ''}"
        if is_instructional:
            return "instruction"
        if any(keyword in content for keyword in ["说明", "备注", "描述", "摘要"]):
            return "description"
        if any(keyword in content for keyword in ["要点", "bullet", "列表"]):
            return "bullet_text"
        return "text"

    def _infer_content_requirement(self, element_type: str, default_text: str | None, shape_name: str, is_instructional: bool) -> str:
        if element_type == "title":
            return "根据需求生成该页标题，保持简洁明确。"
        if default_text and is_instructional:
            return self._clean_instruction_text(default_text)
        if default_text:
            return f"参考模板默认文案风格生成该文本，可改写但不要照抄：{self._clean_instruction_text(default_text)}"
        return f"根据页面用途为 {shape_name} 生成合适内容。"

    def _is_instructional_text(self, shape_name: str, default_text: str | None) -> bool:
        content = f"{shape_name} {default_text or ''}"
        instruction_keywords = [
            "填写说明", "示例", "画法要求", "请填写",
            "可在后续章节", "应体现", "需重点说明",
            "正文", "磅", "行间距",
        ]
        return any(keyword in content for keyword in instruction_keywords)

    def _clean_instruction_text(self, text: str) -> str:
        cleaned = re.sub(r"（[^（）]*?(微软雅黑|宋体|黑体|磅|行距|间距)[^（）]*?）", "", text)
        cleaned = re.sub(r"\([^()]*?(微软雅黑|宋体|黑体|磅|行距|间距)[^()]*?\)", "", cleaned)
        cleaned = re.sub(r"^[【\[]?填写说明[】\]]?", "", cleaned)
        cleaned = re.sub(r"\s+", " ", cleaned)
        return cleaned.strip(" ：:")

    def _local_name(self, tag: str) -> str:
        return tag.rsplit("}", 1)[-1]
