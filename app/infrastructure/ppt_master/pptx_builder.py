"""PPT 构建器。

从 pptx_gen/ppt/builder.py 移植，去掉图片/Mermaid 相关逻辑。
读取模板规则 + 结构化生成结果，直接在模板 PPTX 上回填文本和表格，
支持文本溢出自动拆页和表格分页。
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import json
import re

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from app.schemas.structured_generation import GeneratedElement, StructuredPageResult


class PPTBuilder:
    """在模板 PPTX 上回填结构化内容，支持自动拆页。"""

    def __init__(self, template_path: str | Path) -> None:
        self.template_path = Path(template_path)

    def build(
        self,
        rules: dict,
        page_results: dict[int, StructuredPageResult],
        output_path: str | Path,
    ) -> Path:
        """构建最终 PPTX。

        Args:
            rules: 模板规则 JSON（TemplateRules.model_dump() 的结果）
            page_results: 页码 → StructuredPageResult 的映射
            output_path: 输出 PPTX 路径
        """
        page_rule_map = {int(page["page_no"]): page for page in rules["pages"]}
        slide_size = (rules.get("template") or {}).get("slide_size") or {}

        presentation = Presentation(self.template_path)
        delete_slide_indices: list[int] = []
        original_slide_count = len(presentation.slides)

        for slide_index in reversed(range(original_slide_count)):
            page_no = slide_index + 1
            result = page_results.get(page_no)
            page_rule = page_rule_map.get(page_no)
            if result is None or page_rule is None or not result.should_generate:
                delete_slide_indices.append(slide_index)
                continue
            expanded_results = self._expand_result_for_overflow(slide_size, page_rule, result)
            slide = presentation.slides[slide_index]
            insert_after_index = slide_index
            duplicated_slides: list[Slide] = []
            for _ in expanded_results[1:]:
                duplicated_slide = self._duplicate_slide_after(presentation, insert_after_index)
                duplicated_slides.append(duplicated_slide)
                insert_after_index += 1
            self._fill_slide(slide, page_rule, expanded_results[0])
            for duplicated_slide, extra_result in zip(duplicated_slides, expanded_results[1:]):
                self._fill_slide(duplicated_slide, page_rule, extra_result)

        for slide_index in sorted(delete_slide_indices, reverse=True):
            self._delete_slide(presentation, slide_index)

        target = Path(output_path)
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(target)
        return target

    def fill_single_slide(
        self,
        presentation: Presentation,
        slide_index: int,
        page_rule: dict,
        result: StructuredPageResult,
    ) -> list[int]:
        """回填单个 slide，返回实际占用的 slide 索引列表（含拆页）。

        与 build() 不同，此方法在已打开的 Presentation 上操作，
        不保存文件，不删除跳过的 slide。用于混合导出场景。
        """
        if not result.should_generate:
            return []
        slide_size = {}
        expanded_results = self._expand_result_for_overflow(slide_size, page_rule, result)
        slide = presentation.slides[slide_index]
        insert_after_index = slide_index
        occupied_indices = [slide_index]
        for _ in expanded_results[1:]:
            duplicated_slide = self._duplicate_slide_after(presentation, insert_after_index)
            insert_after_index += 1
            occupied_indices.append(insert_after_index)
        self._fill_slide(slide, page_rule, expanded_results[0])
        for idx, extra_result in enumerate(expanded_results[1:], start=1):
            duplicated_slide = presentation.slides[occupied_indices[idx]]
            self._fill_slide(duplicated_slide, page_rule, extra_result)
        return occupied_indices

    def _fill_slide(self, slide: Slide, page_rule: dict, result: StructuredPageResult) -> None:
        rule_element_map = {item["id"]: item for item in page_rule.get("elements", [])}
        result_element_map = {item.id: item for item in result.elements}
        shape_map = {shape.shape_id: shape for shape in slide.shapes}
        for element_id, rule in rule_element_map.items():
            shape = shape_map.get(int(rule["shape_id"]))
            result_element = result_element_map.get(element_id)
            if rule["type"] in {"title", "text"}:
                if shape is not None and hasattr(shape, "text_frame"):
                    self._set_text(shape, (result_element.content if result_element is not None else "") or "")
            elif rule["type"] == "table":
                if shape is not None and getattr(shape, "has_table", False):
                    if result_element is None or not ((result_element.headers or []) or (result_element.rows or [])):
                        self._remove_shape(shape)
                        continue
                    self._fill_table(shape, rule, result_element)

    def _expand_result_for_overflow(
        self,
        slide_size: dict,
        page_rule: dict,
        result: StructuredPageResult,
    ) -> list[StructuredPageResult]:
        """检测内容溢出并拆页。"""
        if page_rule.get("page_purpose") != "text":
            return [result]
        if any(item.get("type") in {"table"} for item in page_rule.get("elements", [])):
            return self._expand_mixed_result_for_overflow(slide_size, page_rule, result)

        result_element_map = {item.id: item for item in result.elements}
        body_rules = [
            item
            for item in page_rule.get("elements", [])
            if item.get("type") == "text"
            and item.get("id") in result_element_map
            and (result_element_map[item["id"]].content or "").strip()
        ]
        if len(body_rules) != 1:
            return [result]

        body_rule = body_rules[0]
        content_requirement = body_rule.get("content_requirement") or ""
        if "可分多页" not in content_requirement and "分多页" not in content_requirement:
            return [result]

        body_element = result_element_map[body_rule["id"]]
        chunks = self._split_text_to_fit(page_rule, body_rule, body_element.content or "")
        if len(chunks) <= 1:
            return [result]

        return [
            self._clone_result_with_text(result, body_element.id, chunk)
            for chunk in chunks
        ]

    def _expand_mixed_result_for_overflow(
        self,
        slide_size: dict,
        page_rule: dict,
        result: StructuredPageResult,
    ) -> list[StructuredPageResult]:
        bands = self._build_layout_bands(page_rule, result)
        if not bands:
            return [result]

        content_start_top = min(band["original_top"] for band in bands)
        slide_bottom = int(slide_size.get("cy") or 0) or max(
            rule.get("bbox", {}).get("y", 0) + rule.get("bbox", {}).get("h", 0)
            for rule in page_rule.get("elements", [])
        )
        slides: list[StructuredPageResult] = []
        current_element_map: dict[str, GeneratedElement] = {}
        current_bottom = 0
        previous_band = None

        for band in bands:
            segments = self._expand_band_segments(band)
            for segment_index, segment in enumerate(segments):
                if segment_index > 0 and current_element_map:
                    slides.append(self._clone_result_with_element_map(result, current_element_map))
                    current_element_map = {}
                    current_bottom = 0
                    previous_band = None
                gap = 0
                if current_element_map and segment_index == 0 and previous_band is not None:
                    gap = max(int(band["original_top"]) - int(previous_band["original_bottom"]), 0)
                proposed_top = content_start_top if not current_element_map else current_bottom + gap
                if current_element_map and proposed_top + int(segment["required_height"]) > slide_bottom:
                    slides.append(self._clone_result_with_element_map(result, current_element_map))
                    current_element_map = {}
                    current_bottom = 0
                    previous_band = None
                    proposed_top = content_start_top
                for item in segment["items"]:
                    current_element_map[item["element"].id] = item["element"]
                current_bottom = proposed_top + int(segment["required_height"])
                previous_band = band

        if current_element_map:
            slides.append(self._clone_result_with_element_map(result, current_element_map))
        return slides or [result]

    def _build_layout_bands(self, page_rule: dict, result: StructuredPageResult) -> list[dict[str, object]]:
        result_element_map = {item.id: item for item in result.elements}
        active_rules = []
        for rule in sorted(
            page_rule.get("elements", []),
            key=lambda item: (item.get("bbox", {}).get("y", 0), item.get("z_order", 0)),
        ):
            if rule.get("type") == "title":
                continue
            result_element = result_element_map.get(rule["id"])
            if not self._has_visible_result_element(rule, result_element):
                continue
            height = self._estimate_element_height(page_rule, rule, result_element)
            span_height = min(int(rule.get("bbox", {}).get("h", 0) or 0), max(height, 1)) or max(height, 1)
            active_rules.append({
                "rule": rule,
                "element": result_element.model_copy(deep=True),
                "height": height,
                "span_end": int(rule["bbox"]["y"]) + span_height,
            })
        if not active_rules:
            return []

        bands: list[dict[str, object]] = []
        for item in active_rules:
            if not bands or int(item["rule"]["bbox"]["y"]) > int(bands[-1]["span_end"]) + 20000:
                bands.append({
                    "original_top": int(item["rule"]["bbox"]["y"]),
                    "original_bottom": int(item["rule"]["bbox"]["y"]) + int(item["rule"]["bbox"]["h"]),
                    "span_end": int(item["span_end"]),
                    "items": [item],
                })
                continue
            band = bands[-1]
            band["items"].append(item)
            band["original_top"] = min(int(band["original_top"]), int(item["rule"]["bbox"]["y"]))
            band["original_bottom"] = max(
                int(band["original_bottom"]),
                int(item["rule"]["bbox"]["y"]) + int(item["rule"]["bbox"]["h"]),
            )
            band["span_end"] = max(int(band["span_end"]), int(item["span_end"]))
        return bands

    def _expand_band_segments(self, band: dict[str, object]) -> list[dict[str, object]]:
        band_items = band["items"]
        table_chunks: dict[str, list[dict[str, object]]] = {}
        segment_count = 1
        for item in band_items:
            rule = item["rule"]
            if rule.get("type") != "table":
                continue
            chunks = self._chunk_table_element(rule, item["element"])
            table_chunks[rule["id"]] = chunks
            segment_count = max(segment_count, len(chunks))

        segments: list[dict[str, object]] = []
        for segment_index in range(segment_count):
            segment_items: list[dict[str, object]] = []
            required_height = 0
            for item in band_items:
                rule = item["rule"]
                if rule.get("type") == "table":
                    chunks = table_chunks.get(rule["id"], [])
                    if segment_index >= len(chunks):
                        continue
                    chunk = chunks[segment_index]
                    segment_items.append({
                        "rule": rule,
                        "element": chunk["element"],
                        "height": int(chunk["height"]),
                    })
                    required_height = max(required_height, int(chunk["height"]))
                    continue
                segment_items.append({
                    "rule": rule,
                    "element": item["element"].model_copy(deep=True),
                    "height": int(item["height"]),
                })
                required_height = max(required_height, int(item["height"]))
            if segment_items:
                segments.append({"items": segment_items, "required_height": required_height})
        return segments or [{"items": [], "required_height": 0}]

    def _chunk_table_element(self, rule: dict, element: GeneratedElement) -> list[dict[str, object]]:
        headers = list(element.headers or [])
        body_rows = list(element.rows or [])
        template_rows = int(((rule.get("table_schema") or {}).get("rows") or 0))
        header_rows = 1 if headers else 0
        body_capacity = max(template_rows - header_rows, 1)
        if not body_rows:
            chunk_element = element.model_copy(deep=True)
            return [{"element": chunk_element, "height": int(rule.get("bbox", {}).get("h", 0) or 0)}]

        chunks: list[dict[str, object]] = []
        for start in range(0, len(body_rows), body_capacity):
            chunk_element = element.model_copy(deep=True)
            chunk_element.headers = headers or None
            chunk_element.rows = [list(row) for row in body_rows[start : start + body_capacity]]
            chunks.append({
                "element": chunk_element,
                "height": int(rule.get("bbox", {}).get("h", 0) or 0),
            })
        return chunks

    def _has_visible_result_element(self, rule: dict, element: GeneratedElement | None) -> bool:
        if element is None:
            return False
        if rule.get("type") in {"title", "text"}:
            return bool((element.content or "").strip())
        if rule.get("type") == "table":
            return bool((element.headers or []) or (element.rows or []))
        return False

    def _clone_result_with_element_map(
        self,
        result: StructuredPageResult,
        element_map: dict[str, GeneratedElement],
    ) -> StructuredPageResult:
        cloned_elements: list[GeneratedElement] = []
        for element in result.elements:
            mapped = element_map.get(element.id)
            if mapped is not None:
                cloned_elements.append(mapped.model_copy(deep=True))
        return StructuredPageResult(
            page_no=result.page_no,
            should_generate=result.should_generate,
            skip_reason=result.skip_reason,
            elements=cloned_elements,
        )

    def _split_text_to_fit(self, page_rule: dict, rule: dict, text: str) -> list[str]:
        chars_per_line, line_slots = self._estimate_text_capacity(page_rule, rule)
        if chars_per_line <= 0 or line_slots <= 0:
            return [text]

        raw_lines = text.splitlines()
        if not raw_lines:
            return [text]

        wrapped_blocks = [self._wrap_line(raw_line, chars_per_line) for raw_line in raw_lines]
        visual_line_count = sum(max(len(block), 1) for block in wrapped_blocks)
        if visual_line_count <= line_slots:
            return [text]

        chunks: list[str] = []
        current_lines: list[str] = []
        current_visual_lines = 0
        for raw_line, block in zip(raw_lines, wrapped_blocks):
            block_visual_lines = max(len(block), 1)
            if block_visual_lines > line_slots:
                if current_lines:
                    chunks.append(self._trim_chunk_boundary_blank_lines(current_lines))
                    current_lines = []
                    current_visual_lines = 0
                for start in range(0, len(block), line_slots):
                    part = block[start : start + line_slots]
                    chunks.append("".join(part))
                continue

            if current_lines and current_visual_lines + block_visual_lines > line_slots:
                chunks.append(self._trim_chunk_boundary_blank_lines(current_lines))
                current_lines = [raw_line]
                current_visual_lines = block_visual_lines
                continue

            current_lines.append(raw_line)
            current_visual_lines += block_visual_lines
        if current_lines:
            chunks.append(self._trim_chunk_boundary_blank_lines(current_lines))
        return [chunk for chunk in chunks if chunk.strip()]

    def _estimate_text_capacity(self, page_rule: dict, rule: dict) -> tuple[int, int]:
        bbox = rule.get("bbox") or {}
        style = rule.get("style") or {}
        margins = style.get("margins") or {}
        width_pt = max(int(bbox.get("w", 0)) - int(margins.get("left", 0)) - int(margins.get("right", 0)), 0) / 12700
        height_pt = max(int(bbox.get("h", 0)) - int(margins.get("top", 0)) - int(margins.get("bottom", 0)), 0) / 12700
        font_pt = max((style.get("font_size") or 1200) / 100, 10)
        font_hint_pt, line_spacing = self._extract_text_layout_hint(page_rule)
        font_pt = max(font_pt, font_hint_pt)
        chars_per_line = max(int(width_pt / (font_pt * 2.3)), 1)
        line_slots = max(int(height_pt / (font_pt * line_spacing * 1.15)), 1)
        return chars_per_line, line_slots

    def _estimate_element_height(self, page_rule: dict, rule: dict, element: GeneratedElement) -> int:
        if rule.get("type") == "table":
            return max(int(rule.get("bbox", {}).get("h", 0) or 0), int(rule.get("bbox", {}).get("h", 0) or 0))
        if rule.get("type") in {"title", "text"}:
            return max(self._estimate_text_height(page_rule, rule, (element.content or "")), int(rule.get("bbox", {}).get("h", 0) or 0))
        return int(rule.get("bbox", {}).get("h", 0) or 0)

    def _estimate_text_height(self, page_rule: dict, rule: dict, text: str) -> int:
        chars_per_line, _ = self._estimate_text_capacity(page_rule, rule)
        style = rule.get("style") or {}
        margins = style.get("margins") or {}
        font_pt = max((style.get("font_size") or 1200) / 100, 10)
        font_hint_pt, line_spacing = self._extract_text_layout_hint(page_rule)
        font_pt = max(font_pt, font_hint_pt)
        raw_lines = text.splitlines() or [text]
        visual_lines = 0
        for raw_line in raw_lines:
            visual_lines += max(len(self._wrap_line(raw_line, chars_per_line)), 1)
        content_height = int(visual_lines * font_pt * line_spacing * 12700)
        return int(margins.get("top", 0)) + int(margins.get("bottom", 0)) + content_height

    def _extract_text_layout_hint(self, page_rule: dict) -> tuple[float, float]:
        font_pt = 0.0
        line_spacing = 1.4
        for element in page_rule.get("elements", []):
            if element.get("type") != "text":
                continue
            default_text = element.get("default_text") or ""
            font_match = re.search(r"(\d+(?:\.\d+)?)\s*磅", default_text)
            spacing_match = re.search(r"(\d+(?:\.\d+)?)\s*倍行间距", default_text)
            if font_match:
                font_pt = max(font_pt, float(font_match.group(1)))
            if spacing_match:
                line_spacing = max(line_spacing, float(spacing_match.group(1)))
        return font_pt, line_spacing

    def _wrap_line(self, text: str, width: int) -> list[str]:
        if not text:
            return [""]
        if len(text) <= width:
            return [text]
        wrapped: list[str] = []
        remaining = text
        while len(remaining) > width:
            split_at = self._find_split_position(remaining, width)
            wrapped.append(remaining[:split_at])
            remaining = remaining[split_at:]
        if remaining:
            wrapped.append(remaining)
        return wrapped

    def _trim_chunk_boundary_blank_lines(self, lines: list[str]) -> str:
        start = 0
        end = len(lines)
        while start < end and not lines[start].strip():
            start += 1
        while end > start and not lines[end - 1].strip():
            end -= 1
        return "\n".join(lines[start:end])

    def _find_split_position(self, text: str, width: int) -> int:
        candidate = min(width, len(text))
        lower_bound = max(width // 2, 1)
        while candidate > lower_bound:
            if re.match(r"[，。；：、,.;:）)]", text[candidate - 1]):
                return candidate
            candidate -= 1
        return min(width, len(text))

    def _clone_result_with_text(self, result: StructuredPageResult, element_id: str, content: str) -> StructuredPageResult:
        cloned_elements: list[GeneratedElement] = []
        for element in result.elements:
            cloned = element.model_copy(deep=True)
            if cloned.id == element_id:
                cloned.content = content
            cloned_elements.append(cloned)
        return StructuredPageResult(
            page_no=result.page_no,
            should_generate=result.should_generate,
            skip_reason=result.skip_reason,
            elements=cloned_elements,
        )

    def _duplicate_slide_after(self, presentation: Presentation, slide_index: int) -> Slide:
        source_slide = presentation.slides[slide_index]
        duplicated_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        for shape in list(duplicated_slide.shapes):
            self._remove_shape(shape)
        for shape in source_slide.shapes:
            cloned_element = deepcopy(shape._element)
            self._detach_placeholder_metadata(cloned_element)
            duplicated_slide.shapes._spTree.insert_element_before(cloned_element, "p:extLst")
        slide_id_list = presentation.slides._sldIdLst
        new_slide_id = slide_id_list[-1]
        del slide_id_list[-1]
        slide_id_list.insert(slide_index + 1, new_slide_id)
        return presentation.slides[slide_index + 1]

    def _detach_placeholder_metadata(self, shape_element) -> None:
        for node in list(shape_element.iter()):
            if node.tag.endswith("}ph"):
                parent = node.getparent()
                if parent is not None:
                    parent.remove(node)

    def _set_text(self, shape, content: str) -> None:
        text_frame = shape.text_frame
        paragraph_alignment = None
        font_name = None
        font_size = None
        bold = None
        italic = None
        if text_frame.paragraphs:
            paragraph_alignment = text_frame.paragraphs[0].alignment
            if text_frame.paragraphs[0].runs:
                font = text_frame.paragraphs[0].runs[0].font
                font_name = font.name
                font_size = font.size
                bold = font.bold
                italic = font.italic
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = content
        paragraph.alignment = paragraph_alignment
        if paragraph.runs:
            font = paragraph.runs[0].font
            font.name = font_name
            font.size = font_size
            font.bold = bold
            font.italic = italic

    def _fill_table(self, shape, rule: dict, element: GeneratedElement) -> None:
        table = shape.table
        all_rows = []
        if element.headers:
            all_rows.append(element.headers)
        if element.rows:
            all_rows.extend(element.rows)
        style = rule.get("style") or {}
        header_style = self._resolve_table_text_style(style.get("header_style"), is_header=True)
        body_style = self._resolve_table_text_style(style.get("body_style"), is_header=False)
        for row_index in range(len(table.rows)):
            for col_index in range(len(table.columns)):
                value = ""
                if row_index < len(all_rows) and col_index < len(all_rows[row_index]):
                    value = str(all_rows[row_index][col_index])
                cell = table.cell(row_index, col_index)
                cell.text = value
                cell_style = header_style if row_index == 0 and element.headers else body_style
                self._apply_table_cell_style(cell, cell_style)

    def _resolve_table_text_style(self, style: dict[str, object] | None, *, is_header: bool) -> dict[str, object]:
        base = {
            "font_name": "微软雅黑",
            "font_size": 1400,
            "bold": is_header,
            "italic": False,
            "alignment": "center" if is_header else "left",
            "vertical_anchor": "middle",
            "margins": {
                "left": 91440 if is_header else 68580,
                "top": 45720 if is_header else 0,
                "right": 91440 if is_header else 68580,
                "bottom": 45720 if is_header else 0,
            },
        }
        if not style:
            return base
        resolved = dict(base)
        resolved["margins"] = dict(base["margins"])
        for key, value in style.items():
            if key == "margins":
                resolved["margins"].update(value or {})
            else:
                resolved[key] = value
        return resolved

    def _apply_table_cell_style(self, cell, style: dict[str, object]) -> None:
        margins = style.get("margins") or {}
        cell.margin_left = int(margins.get("left", 68580))
        cell.margin_top = int(margins.get("top", 0))
        cell.margin_right = int(margins.get("right", 68580))
        cell.margin_bottom = int(margins.get("bottom", 0))
        cell.vertical_anchor = self._to_vertical_anchor(style.get("vertical_anchor"))
        text_frame = cell.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        paragraph.alignment = self._to_paragraph_alignment(style.get("alignment"))
        if paragraph.runs:
            font = paragraph.runs[0].font
            font.name = style.get("font_name") or "微软雅黑"
            font.size = Pt((style.get("font_size") or 1400) / 100)
            font.bold = bool(style.get("bold", False))
            font.italic = bool(style.get("italic", False))

    def _to_paragraph_alignment(self, alignment: str | None):
        mapping = {
            "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
            "distribute": PP_ALIGN.DISTRIBUTE,
        }
        return mapping.get((alignment or "left").lower(), PP_ALIGN.LEFT)

    def _to_vertical_anchor(self, anchor: str | None):
        mapping = {
            "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        return mapping.get((anchor or "middle").lower(), MSO_ANCHOR.MIDDLE)

    def _remove_shape(self, shape) -> None:
        element = shape._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    def _delete_slide(self, presentation: Presentation, slide_index: int) -> None:
        slide_id_list = presentation.slides._sldIdLst
        slide_id = slide_id_list[slide_index]
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        del slide_id_list[slide_index]
