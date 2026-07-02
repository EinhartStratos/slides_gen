"""混合导出器。

将 SVG 页面（通过 convert_svg_to_slide_shapes 转为 DrawingML 可编辑形状）
和结构化页面（通过 PPTBuilder 回填到模板原生元素）整合到同一个 PPTX 中。

整合流程：
1. 用 python-pptx 打开模板 PPTX 副本
2. 结构化页面 → PPTBuilder.fill_single_slide() 回填对应 slide
3. SVG 页面 → convert_svg_to_slide_shapes() 生成 DrawingML XML → 注入到对应 slide
4. 删除跳过的页面
5. 重排 slide 顺序
6. 保存最终 PPTX
"""

from __future__ import annotations

from copy import deepcopy
from io import BytesIO
import logging
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation

from app.core.config import Settings
from app.infrastructure.ppt_master.pptx_builder import PPTBuilder
from app.schemas.structured_generation import StructuredPageResult


logger = logging.getLogger(__name__)

# DrawingML 命名空间
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _qn(tag: str) -> str:
    """构建带命名空间的标签。"""
    prefix, local = tag.split(":")
    ns_map = {"p": NS_P, "a": NS_A, "r": NS_R}
    return f"{{{ns_map[prefix]}}}{local}"


class HybridPptxExporter:
    """混合导出器：整合 SVG slide + 结构化 slide 到一个 PPTX。"""

    def __init__(
        self,
        settings: Settings,
        pptx_builder: PPTBuilder | None = None,
    ) -> None:
        self.settings = settings
        self._svg_convert_func = None

    def _get_svg_convert_func(self):
        """延迟加载 convert_svg_to_slide_shapes 函数。"""
        if self._svg_convert_func is not None:
            return self._svg_convert_func
        scripts_dir = self.settings.ppt_master_scripts_dir
        if str(scripts_dir) not in sys.path:
            sys.path.insert(0, str(scripts_dir))
        from svg_to_pptx.drawingml_converter import convert_svg_to_slide_shapes
        self._svg_convert_func = convert_svg_to_slide_shapes
        return self._svg_convert_func

    def export(
        self,
        template_pptx_path: Path,
        template_rules: dict,
        svg_pages: dict[int, Path],
        structured_pages: dict[int, StructuredPageResult],
        skipped_pages: set[int],
        output_path: Path,
    ) -> Path:
        """整合两种来源的页面，导出最终 PPTX。

        Args:
            template_pptx_path: 模板 PPTX 路径
            template_rules: 模板规则 JSON dict
            svg_pages: 页码 → SVG 文件路径
            structured_pages: 页码 → StructuredPageResult
            skipped_pages: 跳过的页码集合
            output_path: 输出 PPTX 路径
        """
        presentation = Presentation(template_pptx_path)
        original_slide_count = len(presentation.slides)
        page_rule_map = {int(page["page_no"]): page for page in template_rules.get("pages", [])}

        # 页码 → 实际占用的 slide 索引列表（结构化页面可能拆页产生多个 slide）
        page_to_slide_indices: dict[int, list[int]] = {}
        # 需要删除的 slide 索引
        delete_indices: set[int] = set()

        for slide_index in range(original_slide_count):
            page_no = slide_index + 1
            if page_no in skipped_pages:
                delete_indices.add(slide_index)
                page_to_slide_indices[page_no] = []
                continue

            if page_no in structured_pages:
                result = structured_pages[page_no]
                if not result.should_generate:
                    delete_indices.add(slide_index)
                    page_to_slide_indices[page_no] = []
                    continue
                page_rule = page_rule_map.get(page_no, {})
                builder = PPTBuilder(template_pptx_path)
                occupied = builder.fill_single_slide(presentation, slide_index, page_rule, result)
                page_to_slide_indices[page_no] = occupied

            elif page_no in svg_pages:
                svg_file = svg_pages[page_no]
                self._inject_svg_slide(presentation, slide_index, svg_file)
                page_to_slide_indices[page_no] = [slide_index]

            else:
                delete_indices.add(slide_index)
                page_to_slide_indices[page_no] = []

        # 删除跳过的 slide
        for slide_index in sorted(delete_indices, reverse=True):
            self._delete_slide(presentation, slide_index)
            # 调整后续页码的 slide 索引映射
            for page_no, indices in page_to_slide_indices.items():
                page_to_slide_indices[page_no] = [
                    idx - 1 if idx > slide_index else idx for idx in indices
                ]

        # 重排 slide 顺序：按页码顺序排列
        self._reorder_slides(presentation, page_to_slide_indices)

        # 保存
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return output_path

    def _inject_svg_slide(
        self,
        presentation: Presentation,
        slide_index: int,
        svg_file: Path,
    ) -> None:
        """将 SVG 转换的 DrawingML 注入到模板 PPTX 的指定 slide。

        保持可编辑形状，不渲染为图片。
        """
        convert_func = self._get_svg_convert_func()
        slide_xml, media_files, rel_entries, _ = convert_func(
            svg_file, slide_num=slide_index + 1, verbose=False,
        )

        # 解析转换后的 XML
        new_slide_root = etree.fromstring(slide_xml.encode("utf-8"))
        new_sp_tree = new_slide_root.find(_qn("p:spTree"))
        if new_sp_tree is None:
            logger.warning("SVG 转换结果中未找到 spTree (slide=%s)", slide_index)
            return

        # 获取目标 slide 的 spTree
        target_slide = presentation.slides[slide_index]
        target_element = target_slide._element
        target_cSld = target_element.find(_qn("p:cSld"))
        if target_cSld is None:
            logger.warning("目标 slide 中未找到 cSld (slide=%s)", slide_index)
            return
        target_sp_tree = target_cSld.find(_qn("p:spTree"))
        if target_sp_tree is None:
            logger.warning("目标 slide 中未找到 spTree (slide=%s)", slide_index)
            return

        # 删除目标 slide 中所有现有 shape（保留 nvGrpSpPr 和 grpSpPr）
        for child in list(target_sp_tree):
            tag = etree.QName(child).localname
            if tag not in ("nvGrpSpPr", "grpSpPr"):
                target_sp_tree.remove(child)

        # 计算 shape ID 偏移量，避免与已有 slide 的 shape ID 冲突
        max_existing_id = self._find_max_shape_id(presentation)
        id_offset = max_existing_id

        # 将新 shape 插入目标 slide（跳过新 spTree 的 nvGrpSpPr 和 grpSpPr）
        for child in new_sp_tree:
            tag = etree.QName(child).localname
            if tag in ("nvGrpSpPr", "grpSpPr"):
                continue
            cloned = deepcopy(child)
            # 偏移所有 shape ID
            self._offset_shape_ids(cloned, id_offset)
            target_sp_tree.append(cloned)

        # 处理 media 文件和 relationships
        if media_files:
            slide_part = target_slide.part
            for filename, data in media_files.items():
                try:
                    image_part = slide_part.get_or_add_image_part(BytesIO(data))
                    logger.debug("已添加图片到 slide %s: %s", slide_index, filename)
                except Exception as exc:
                    logger.warning("添加图片失败 (slide=%s, file=%s): %s", slide_index, filename, exc)

        logger.info("SVG slide 注入完成 (slide=%s, svg=%s)", slide_index, svg_file.name)

    def _find_max_shape_id(self, presentation: Presentation) -> int:
        """找到所有 slide 中最大的 shape ID。"""
        max_id = 100
        for slide in presentation.slides:
            for shape in slide.shapes:
                try:
                    shape_id = shape.shape_id
                    if shape_id > max_id:
                        max_id = shape_id
                except Exception:
                    pass
        return max_id + 1

    def _offset_shape_ids(self, element, offset: int) -> None:
        """递归偏移 XML 元素中的所有 shape ID。"""
        tag = etree.QName(element).localname
        if tag == "cNvPr" and "id" in element.attrib:
            try:
                old_id = int(element.attrib["id"])
                element.attrib["id"] = str(old_id + offset)
            except (ValueError, TypeError):
                pass
        for child in element:
            self._offset_shape_ids(child, offset)

    def _delete_slide(self, presentation: Presentation, slide_index: int) -> None:
        """删除指定索引的 slide。"""
        slide_id_list = presentation.slides._sldIdLst
        slide_id = slide_id_list[slide_index]
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        del slide_id_list[slide_index]

    def _reorder_slides(
        self,
        presentation: Presentation,
        page_to_slide_indices: dict[int, list[int]],
    ) -> None:
        """按页码顺序重排 slide。"""
        sorted_pages = sorted(page_to_slide_indices.keys())
        new_order: list[int] = []
        for page_no in sorted_pages:
            new_order.extend(page_to_slide_indices[page_no])

        if not new_order:
            return

        slide_id_list = presentation.slides._sldIdLst
        # 收集所有 sldId 元素
        sld_ids = [slide_id_list[i] for i in range(len(slide_id_list))]

        # 按新顺序重新排列
        for i, target_index in enumerate(new_order):
            if target_index < len(sld_ids):
                sld_id = sld_ids[target_index]
                # 先移除再插入到正确位置
                slide_id_list.remove(sld_id)
                slide_id_list.insert(i, sld_id)
